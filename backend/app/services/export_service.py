import io
import logging
from typing import Optional

from app.core.database import get_db
from app.schemas.export import ExportFormat
from fastapi import HTTPException
from fastapi.responses import StreamingResponse

logger = logging.getLogger(__name__)


class ExportService:
    async def export(
        self,
        entity_type: str,
        entity_id: str,
        format: ExportFormat,
        user_id: str,
        title: Optional[str] = None,
        include_citations: bool = True,
    ) -> StreamingResponse:
        db = get_db()
        entity = await self._get_entity(db, entity_type, entity_id, user_id)
        if not entity:
            raise HTTPException(status_code=404, detail=f"{entity_type.capitalize()} not found")

        content_text = self._extract_content(entity, entity_type)
        doc_title = title or entity.get("title", "Atlas AI Export")
        citations = entity.get("citations", []) if include_citations else []

        if format == ExportFormat.MARKDOWN:
            return self._export_markdown(doc_title, content_text, citations)
        elif format == ExportFormat.PDF:
            return await self._export_pdf(doc_title, content_text, citations)
        elif format == ExportFormat.DOCX:
            return await self._export_docx(doc_title, content_text, citations)
        elif format == ExportFormat.CSV:
            return self._export_csv(entity, entity_type)
        elif format == ExportFormat.PPTX:
            return await self._export_pptx(doc_title, content_text)
        else:
            raise HTTPException(status_code=400, detail="Unsupported format")

    async def _get_entity(self, db, entity_type: str, entity_id: str, user_id: str):
        collection_map = {"note": "notes", "artifact": "artifacts", "chat": "chats"}
        col = collection_map.get(entity_type)
        if not col:
            return None
        doc = await db[col].find_one({"_id": entity_id, "user_id": user_id})
        return doc

    def _extract_content(self, entity: dict, entity_type: str) -> str:
        if entity_type == "note":
            return entity.get("content_text") or entity.get("content_html") or ""
        elif entity_type == "artifact":
            content = entity.get("content")
            if isinstance(content, str):
                return content
            elif isinstance(content, dict):
                return "\n".join(f"{k}: {v}" for k, v in content.items())
            elif isinstance(content, list):
                return "\n".join(str(item) for item in content)
            return str(content)
        elif entity_type == "chat":
            messages = entity.get("messages", [])
            parts = []
            for msg in messages:
                role = msg.get("role", "").upper()
                content = msg.get("content", "")
                parts.append(f"**{role}:** {content}")
            return "\n\n".join(parts)
        return ""

    def _export_markdown(self, title: str, content: str, citations: list) -> StreamingResponse:
        md = f"# {title}\n\n{content}"
        if citations:
            md += "\n\n## Citations\n"
            for i, c in enumerate(citations, 1):
                md += f"\n{i}. **{c.get('source_name', 'Unknown')}**"
                if c.get("page_number"):
                    md += f" (p. {c['page_number']})"
                if c.get("text_excerpt"):
                    md += f"\n   > {c['text_excerpt'][:200]}"
        buf = io.BytesIO(md.encode("utf-8"))
        return StreamingResponse(
            buf,
            media_type="text/markdown",
            headers={"Content-Disposition": f'attachment; filename="{title}.md"'},
        )

    async def _export_pdf(self, title: str, content: str, citations: list) -> StreamingResponse:
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.units import inch

            buf = io.BytesIO()
            doc = SimpleDocTemplate(buf, pagesize=letter)
            styles = getSampleStyleSheet()
            story = [
                Paragraph(title, styles["Title"]),
                Spacer(1, 0.3 * inch),
                Paragraph(content.replace("\n", "<br/>"), styles["Normal"]),
            ]
            if citations:
                story.append(Spacer(1, 0.3 * inch))
                story.append(Paragraph("Citations", styles["Heading2"]))
                for i, c in enumerate(citations, 1):
                    cit_text = f"{i}. {c.get('source_name', 'Unknown')}"
                    if c.get("page_number"):
                        cit_text += f" (p. {c['page_number']})"
                    story.append(Paragraph(cit_text, styles["Normal"]))
            doc.build(story)
            buf.seek(0)
            return StreamingResponse(
                buf,
                media_type="application/pdf",
                headers={"Content-Disposition": f'attachment; filename="{title}.pdf"'},
            )
        except ImportError:
            raise HTTPException(status_code=501, detail="PDF export requires reportlab: pip install reportlab")

    async def _export_docx(self, title: str, content: str, citations: list) -> StreamingResponse:
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument()
            doc.add_heading(title, 0)
            for para in content.split("\n\n"):
                doc.add_paragraph(para)
            if citations:
                doc.add_heading("Citations", level=1)
                for i, c in enumerate(citations, 1):
                    text = f"{i}. {c.get('source_name', 'Unknown')}"
                    if c.get("page_number"):
                        text += f" (p. {c['page_number']})"
                    doc.add_paragraph(text)
            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)
            return StreamingResponse(
                buf,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={"Content-Disposition": f'attachment; filename="{title}.docx"'},
            )
        except ImportError:
            raise HTTPException(status_code=501, detail="DOCX export requires python-docx")

    def _export_csv(self, entity: dict, entity_type: str) -> StreamingResponse:
        import csv
        buf = io.StringIO()
        writer = csv.writer(buf)
        if entity_type == "artifact":
            content = entity.get("content")
            if isinstance(content, list) and content and isinstance(content[0], dict):
                writer.writerow(content[0].keys())
                for row in content:
                    writer.writerow(row.values())
            else:
                writer.writerow(["content"])
                writer.writerow([str(content)])
        else:
            writer.writerow(["field", "value"])
            for k, v in entity.items():
                if k not in ("_id", "hashed_password"):
                    writer.writerow([k, str(v)])
        buf.seek(0)
        return StreamingResponse(
            io.BytesIO(buf.getvalue().encode("utf-8")),
            media_type="text/csv",
            headers={"Content-Disposition": f'attachment; filename="export.csv"'},
        )

    async def _export_pptx(self, title: str, content: str) -> StreamingResponse:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            tf.text = content[:800]
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            return StreamingResponse(
                buf,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f'attachment; filename="{title}.pptx"'},
            )
        except ImportError:
            raise HTTPException(status_code=501, detail="PPTX export requires python-pptx")


export_service = ExportService()
