import io
import logging
from typing import Optional
from pathlib import Path

logger = logging.getLogger(__name__)


async def extract_pages_from_pdf(file_path: str) -> tuple[list[str], int]:
    try:
        import fitz
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        page_count = len(doc)
        doc.close()
        return pages, page_count
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise


async def extract_text_from_pdf(file_path: str) -> tuple[str, int]:
    pages, page_count = await extract_pages_from_pdf(file_path)
    return "\n".join(pages), page_count


async def extract_text_from_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        return "\n".join(paragraphs)
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        raise


async def extract_text_from_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()
    except Exception as e:
        logger.error(f"TXT extraction error: {e}")
        raise


async def extract_text_from_csv(file_path: str) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    except Exception as e:
        logger.error(f"CSV extraction error: {e}")
        raise


async def extract_text_from_xlsx(file_path: str) -> str:
    try:
        import pandas as pd
        excel_file = pd.ExcelFile(file_path)
        parts = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            parts.append(f"=== Sheet: {sheet_name} ===\n{df.to_string(index=False)}")
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"XLSX extraction error: {e}")
        raise


async def extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                parts.append(f"=== Slide {i} ===\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise


async def extract_text_from_url(url: str) -> str:
    try:
        import httpx
        from bs4 import BeautifulSoup
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            response = await client.get(url)
            response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"URL extraction error: {e}")
        raise
